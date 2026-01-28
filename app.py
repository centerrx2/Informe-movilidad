import io
import os
import zipfile
import tempfile
from typing import Dict, Any, List

import numpy as np
import streamlit as st
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
import xml.etree.ElementTree as ET
from scipy.spatial.transform import Rotation as R
from scipy.signal import find_peaks
from reportlab.lib.utils import ImageReader
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image
import matplotlib.pyplot as plt
from reportlab.lib import colors
from reportlab.lib.units import cm
from reportlab.lib.pagesizes import A4
from reportlab.platypus import PageBreak

from reportlab.lib.styles import ParagraphStyle
from reportlab.platypus import  KeepTogether

from datetime import datetime
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

IMAGE_MAP = {
    "ASLR": os.path.join(BASE_DIR, "fotos", "aslr_sin_fondo.png"),
    "AKE": os.path.join(BASE_DIR, "fotos", "ake_sin_fondo.png"),
    "Rotación de cadera": os.path.join(BASE_DIR, "fotos", "cadera_sin_fondo.png"),
    "Rotación torácica": os.path.join(BASE_DIR, "fotos", "rotacion_marcado.png"),
    "Thomas": os.path.join(BASE_DIR, "fotos", "thomas_sin_fondo.png"),
    "Lunge": os.path.join(BASE_DIR, "fotos", "lunge_sin_fondo_2.png"),
    "Plantar": os.path.join(BASE_DIR, "fotos", "plantar_sin_fondo.png"),
    "OHS": os.path.join(BASE_DIR, "fotos", "ohs_sin_fondo_2.png"),
    "OHSf": os.path.join(BASE_DIR, "fotos", "ohs_frontal_sin_fondo.png"),
    "Step Down": os.path.join(BASE_DIR, "fotos", "sdown_sin_fondo.png"),
}


# ===================================================================
#  1. MAPAS DE JOINTS Y SEGMENTOS (basado en tus .mvnx reales)
# ===================================================================

JOINT_ORDER = [
    "jL5S1",
    "jL4L3",
    "jL1T12",
    "jT9T8",
    "jT1C7",
    "jC1Head",
    "jRightT4Shoulder",
    "jRightShoulder",
    "jRightElbow",
    "jRightWrist",
    "jLeftT4Shoulder",
    "jLeftShoulder",
    "jLeftElbow",
    "jLeftWrist",
    "jRightHip",      # index 14
    "jRightKnee",     # index 15
    "jRightAnkle",    # index 16
    "jRightBallFoot",
    "jLeftHip",       # index 18
    "jLeftKnee",      # index 19
    "jLeftAnkle",     # index 20
    "jLeftBallFoot",
]

JOINT_MAP = {
    "RightHip": JOINT_ORDER.index("jRightHip"),
    "RightKnee": JOINT_ORDER.index("jRightKnee"),
    "RightAnkle": JOINT_ORDER.index("jRightAnkle"),
    "LeftHip": JOINT_ORDER.index("jLeftHip"),
    "LeftKnee": JOINT_ORDER.index("jLeftKnee"),
    "LeftAnkle": JOINT_ORDER.index("jLeftAnkle"),
}

SEGMENT_LABELS = {
    1: "Pelvis",
    2: "L5",
    3: "L3",
    4: "T12",
    5: "T8",
    6: "Neck",
    7: "Head",
    8: "RightShoulder",
    9: "RightUpperArm",
    10: "RightForeArm",
    11: "RightHand",
    12: "LeftShoulder",
    13: "LeftUpperArm",
    14: "LeftForeArm",
    15: "LeftHand",
    16: "RightUpperLeg",
    17: "RightLowerLeg",
    18: "RightFoot",
    19: "RightToe",
    20: "LeftUpperLeg",
    21: "LeftLowerLeg",
    22: "LeftFoot",
    23: "LeftToe"
}
LOGO_PORTADA = os.path.join(BASE_DIR, "fotos", "logo.png")

# ===================================================================
#  1.1 CONFIGURACIÓN XSENS (AWINDA vs LINK)
# ===================================================================

XSENS_CONFIG = {
    "Awinda": {
        "ROT_CADERA": {
            "R": 3,    # right lower leg
            "L": 6     # left lower leg
        },
        "THOMAS": {
            "R": 2,    # right upper leg
            "L": 5     # left upper leg
        }
    },
    "Link (traje)": {
        "ROT_CADERA": {
            "R": 10,
            "L": 13
        },
        "THOMAS": {
            "R": 9,
            "L": 12
        }
    }
}


# ===================================================================
#  2. FUNCIONES DE UTILIDAD
# ===================================================================
def quat2eul_matlab(q, order="ZYX"):
    """
    Reproduce MATLAB quat2eul(q, order)
    MVNX: q = [qw qx qy qz]
    Devuelve ángulos en grados
    """
    qw, qx, qy, qz = q.T

    # -------- ZYX (yaw–pitch–roll) --------
    if order == "ZYX":
        z = np.arctan2(
            2*(qw*qz + qx*qy),
            1 - 2*(qy*qy + qz*qz)
        )

        siny = 2*(qw*qy - qz*qx)
        siny = np.clip(siny, -1.0, 1.0)
        y = np.arcsin(siny)

        x = np.arctan2(
            2*(qw*qx + qy*qz),
            1 - 2*(qx*qx + qy*qy)
        )

        return np.rad2deg(np.column_stack([z, y, x]))

    # -------- YZX (MATLAB Thomas test) --------
    elif order == "YZX":
        y = np.arctan2(
            2*(qw*qy + qx*qz),
            1 - 2*(qy*qy + qz*qz)
        )

        sinz = 2*(qw*qz - qx*qy)
        sinz = np.clip(sinz, -1.0, 1.0)
        z = np.arcsin(sinz)

        x = np.arctan2(
            2*(qw*qx + qy*qz),
            1 - 2*(qx*qx + qz*qz)
        )

        return np.rad2deg(np.column_stack([y, z, x]))

    else:
        raise ValueError(f"Orden no soportado: {order}")
    


def cut_at_nth_max(signal, n_peaks=4, min_distance=200, min_prominence=5):
    """
    Corta en el frame del n-ésimo pico de FLEXIÓN (1 pico = 1 repetición)
    """

    signal = np.asarray(signal)

    # 🔑 flexión como pico positivo
    flex = -signal

    peaks, _ = find_peaks(
        flex,
        distance=min_distance,
        prominence=min_prominence
    )

    if len(peaks) < n_peaks:
        return len(signal)

    return int(peaks[n_peaks - 1])



def ake_mins_until_cut(
    signal,
    cut_idx,
    min_prominence=10,
    min_distance=150
):
    """
    Detecta mínimos SOLO hasta el frame cut_idx
    """

    signal = np.asarray(signal[:cut_idx])
    inv = -signal

    mins, _ = find_peaks(
        inv,
        prominence=min_prominence,
        distance=min_distance
    )

    return [(signal[p], int(p)) for p in mins]


def mvnx_to_scipy_quat(q):
    """
    MVNX:  [qw qx qy qz]
    SciPy: [qx qy qz qw]
    """
    qw = q[:, 0]
    qx = q[:, 1]
    qy = q[:, 2]
    qz = q[:, 3]
    return np.column_stack([qx, qy, qz, qw])



def get_fecha_from_mvnx(mvnx_path: str):
    NS = "{http://www.xsens.com/mvn/mvnx}"

    root = ET.parse(mvnx_path).getroot()

    subject = root.find(NS + "subject")
    if subject is None:
        return None

    # ✅ 1) usar milisegundos desde Epoch (mejor opción)
    msecs = subject.attrib.get("recDateMSecsSinceEpoch")
    if msecs is not None:
        ts = int(msecs) / 1000.0
        return datetime.fromtimestamp(ts)

    # ✅ 2) fallback: usar recDate (texto)
    rec_txt = subject.attrib.get("recDate")
    if rec_txt is not None:
        # Ejemplo: "Thu Nov 27 12:13:17.238 2025"
        try:
            return datetime.strptime(rec_txt, "%a %b %d %H:%M:%S.%f %Y")
        except:
            pass

    return None


def is_valid_mvnx(path):
    try:
        with open(path, "rb") as f:
            head = f.read(200).lstrip()
        return head.startswith(b"<?xml") or head.startswith(b"<mvnx")
    except:
        return False


def set_cell_center(cell):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER


def set_cell_bg(cell, rgb):
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)

def delete_slide(prs, slide_index):
    """
    Elimina una diapositiva por índice en python-pptx
    """
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)

    prs.part.drop_rel(slides[slide_index].rId)
    xml_slides.remove(slides[slide_index])

def set_metric_text_12(cell, text):
    p = cell.text_frame.paragraphs[0]

    if p.runs:
        run = p.runs[0]
        run.text = str(text)
    else:
        run = p.add_run()
        run.text = str(text)

    run.font.size = Pt(12)




# ===================================================================
#  3. CARGA REAL DEL .MVNX → ADAPTADA A TUS ARCHIVOS
# ===================================================================

def load_mvnx_struct(path: str) -> Dict[str, Any]:
    NS = "{http://www.xsens.com/mvn/mvnx}"
    root = ET.parse(path).getroot()

    subject = root.find(NS + "subject")
    frames = subject.find(NS + "frames")

    orientations = []
    joint_angles = []
    sensor_orientations = []
    joint_angles_ergo = []


    for f in frames.findall(NS + "frame"):
        # orientación → 92 floats → 23 segmentos x 4
        o = f.find(NS + "orientation")
        if o is not None:
            vals = np.array(list(map(float, o.text.split()))).reshape(23, 4)
            orientations.append(vals)

        # jointAngle → 66 floats → 22 joints x 3
        ja = f.find(NS + "jointAngle")
        if ja is not None and ja.text:
            vals = np.array(list(map(float, ja.text.split()))).reshape(22, 3)
            joint_angles.append(vals)



        jae = f.find(NS + "jointAngleErgo")
        if jae is not None:
            vals = np.array(list(map(float, jae.text.split())))
            if len(vals) % 3 == 0:
                vals = vals.reshape(-1, 3)  # ⚠️ NO asumimos 22
                joint_angles_ergo.append(vals)
            else:
                joint_angles_ergo.append(None)
        else:
            joint_angles_ergo.append(None)


                # sensorOrientation → sensores x 4
        so = f.find(NS + "sensorOrientation")
        if so is not None:
            vals = np.array(list(map(float, so.text.split()))).reshape(-1, 4)
            sensor_orientations.append(vals)


    orientations = np.array(orientations)       # shape (N, 23, 4)
    joint_angles = np.array(joint_angles)       # shape (N, 22, 3)
    joint_angles_ergo = np.array(joint_angles_ergo, dtype=object)
    sensor_orientations = np.array(sensor_orientations)  # (N_frames, N_sensors, 4)



    # Construimos en formato MATLAB
    segmentData = {
        sid: {"orientation": orientations[:, sid - 1, :]}
        for sid in SEGMENT_LABELS
    }

    jointData = {
        idx: {"jointAngle": joint_angles[:, idx, :]}
        for idx in range(joint_angles.shape[1])
    }

    sensorData = {
    idx: {"orientation": sensor_orientations[:, idx, :]}
    for idx in range(sensor_orientations.shape[1])
    }


    return {
        "segmentData": segmentData,
        "jointData": jointData,
        "sensorData": sensorData,
        "jointAngleErgo": joint_angles_ergo
    }

# ===================================================================
#  4. CÁLCULOS 
# ===================================================================


def compute_aslr(tree_aslr, start_frame=100, plot_right=True):

    # ===================== DERECHA (segmento 16) =====================
    R_seg_q = tree_aslr["segmentData"][16]["orientation"]  # (N,4) -> [qw qx qy qz]
    R_eul = quat2eul_matlab(R_seg_q, order="ZYX")          # (N,3) -> [Z, Y, X] en grados

    R_signal = R_eul[:, 1]  # columna Y (pitch)

    # ✅ máximo a partir del frame 100
    R_signal_cut = R_signal[start_frame:]
    max_R = float(np.max(R_signal_cut))
    idx_R = int(np.argmax(R_signal_cut) + start_frame)  # frame real


    if max_R > 0:
        R_angle = 90 + max_R
    else:
        R_angle = 90 - abs(max_R)

    # ===================== IZQUIERDA (segmento 20) =====================
    L_seg_q = tree_aslr["segmentData"][20]["orientation"]  # (N,4) -> [qw qx qy qz]
    L_eul = quat2eul_matlab(L_seg_q, order="ZYX")          # (N,3) -> [Z, Y, X] en grados

    L_signal = L_eul[:, 1]  # columna Y (pitch)

    L_signal_cut = L_signal[start_frame:]
    max_L = float(np.max(L_signal_cut))
    idx_L = int(np.argmax(L_signal_cut) + start_frame)


    if max_L > 0:
        L_angle = 90 + max_L
    else:
        L_angle = 90 - abs(max_L)

    return float(R_angle), float(L_angle)



def compute_ake(tree_ake):

    knee_R = tree_ake["jointData"][JOINT_MAP["RightKnee"]]["jointAngle"][:, 2]
    knee_L = tree_ake["jointData"][JOINT_MAP["LeftKnee"]]["jointAngle"][:, 2]

    # -------- DERECHA --------
    cut_R = cut_at_nth_max(knee_R, n_peaks=4)
    mins_R = ake_mins_until_cut(knee_R, cut_R)

    # -------- IZQUIERDA --------
    cut_L = cut_at_nth_max(knee_L, n_peaks=4)
    mins_L = ake_mins_until_cut(knee_L, cut_L)

    if len(mins_R) == 0 or len(mins_L) == 0:
        return None, None, None, None

    R_val, R_frame = min(mins_R, key=lambda x: x[0])
    L_val, L_frame = min(mins_L, key=lambda x: x[0])



    return abs(float(R_val)), abs(float(L_val)), R_frame, L_frame

def compute_rot_cadera(tree, xsens_mode):
    cfg = XSENS_CONFIG[xsens_mode]["ROT_CADERA"]

    # DERECHA
    R_seg = tree["sensorData"][cfg["R"]]["orientation"]
    eu_R = quat2eul_matlab(R_seg)

    R_offset = eu_R[0, 0]
    R_internal = abs(np.min(eu_R[:, 0]) + R_offset)
    R_external = abs(np.max(eu_R[:, 0]) - R_offset)

    # IZQUIERDA
    L_seg = tree["sensorData"][cfg["L"]]["orientation"]
    eu_L = quat2eul_matlab(L_seg)

    L_offset = eu_L[0, 0]
    L_internal = abs(np.max(eu_L[:, 0]) - L_offset)
    L_external = abs(np.min(eu_L[:, 0]) + L_offset)

    return R_internal, R_external, L_internal, L_external



def compute_rot_toracica(tree_rot):
    joint_idx = 4 - 1
    col_idx = 2 - 1

    signal = []

    for fr in tree_rot["jointAngleErgo"]:
        if fr is not None and fr.shape[0] > joint_idx:
            signal.append(fr[joint_idx, col_idx])
        else:
            signal.append(np.nan)

    signal = np.array(signal)

    R_torso = abs(float(np.nanmin(signal)))
    L_torso = abs(float(np.nanmax(signal)))

    return R_torso, L_torso


def compute_thomas(tree_d, tree_i, xsens_mode, start_frame=0):
    cfg = XSENS_CONFIG[xsens_mode]["THOMAS"]

    # ===================== DERECHA =====================
    q_R = tree_d["sensorData"][cfg["R"]]["orientation"]
    eu_R = quat2eul_matlab(q_R, order="ZYX")
    signal_R = eu_R[:, 1]

    idx_R = int(np.argmax(signal_R[start_frame:]) + start_frame)
    R_thigh = float(signal_R[idx_R])

    knee_R_signal = tree_d["jointData"][JOINT_MAP["RightKnee"]]["jointAngle"][:, 2]
    R_knee = float(knee_R_signal[idx_R])

    # ===================== IZQUIERDA =====================
    q_L = tree_i["sensorData"][cfg["L"]]["orientation"]
    eu_L = quat2eul_matlab(q_L, order="ZYX")
    signal_L = eu_L[:, 1]

    idx_L = int(np.argmax(signal_L[start_frame:]) + start_frame)
    L_thigh = float(signal_L[idx_L])

    knee_L_signal = tree_i["jointData"][JOINT_MAP["LeftKnee"]]["jointAngle"][:, 2]
    L_knee = float(knee_L_signal[idx_L])

    return R_thigh, R_knee, L_thigh, L_knee





def compute_lunge(tree_lunge):
    R = tree_lunge["jointData"][JOINT_MAP["RightAnkle"]]["jointAngle"][:,2]
    L = tree_lunge["jointData"][JOINT_MAP["LeftAnkle"]]["jointAngle"][:,2]
    return abs(float(np.max(R))), abs(float(np.max(L)))

 
def compute_plantarflex(tree_plantar):
    R = tree_plantar["jointData"][JOINT_MAP["RightAnkle"]]["jointAngle"][:, 2]
    L = tree_plantar["jointData"][JOINT_MAP["LeftAnkle"]]["jointAngle"][:, 2]
    return abs(float(np.min(R))), abs(float(np.min(L)))


def compute_ohs(tree_ohs):
    # ================== TRONCO Y PELVIS ==================
    stern = tree_ohs["segmentData"][5]["orientation"]
    pelvis = tree_ohs["segmentData"][1]["orientation"]

    eu_s = quat2eul_matlab(stern, "ZYX")
    eu_p = quat2eul_matlab(pelvis, "ZYX")

    # Inclinación anterior del tronco
    sternum_max = float(np.max(eu_s[:, 1]))

    # Inclinación lateral de la pelvis
    pelvis_max = float(np.max(eu_p[:, 2]))
    pelvis_min = float(np.min(eu_p[:, 2]))

    if pelvis_max <= 5 and pelvis_min >= -5:
        pelvis_lat = "NO"
    else:
        pelvis_lat = (
            pelvis_max
            if abs(pelvis_max) >= abs(pelvis_min)
            else pelvis_min
        )

    # ================== CADERA ==================
    hips_R = tree_ohs["jointData"][JOINT_MAP["RightHip"]]["jointAngle"][:, 2]
    hips_L = tree_ohs["jointData"][JOINT_MAP["LeftHip"]]["jointAngle"][:, 2]

    hip_R = float(np.max(hips_R))
    hip_L = float(np.max(hips_L))

    # ================== RODILLA ==================
    knees_R = tree_ohs["jointData"][JOINT_MAP["RightKnee"]]["jointAngle"][:, 2]
    knees_L = tree_ohs["jointData"][JOINT_MAP["LeftKnee"]]["jointAngle"][:, 2]

    knee_R = float(np.max(knees_R))
    knee_L = float(np.max(knees_L))

    # -------- Valgo rodilla (SOLO MÁXIMO) --------
    abd_R = tree_ohs["jointData"][JOINT_MAP["RightKnee"]]["jointAngle"][:, 0]
    abd_L = tree_ohs["jointData"][JOINT_MAP["LeftKnee"]]["jointAngle"][:, 0]

    abd_R_max = float(np.max(abd_R))
    abd_L_max = float(np.max(abd_L))

    valgo_R = abd_R_max if abd_R_max > 5 else "NO"
    valgo_L = abd_L_max if abd_L_max > 5 else "NO"

    # ================== TOBILLO ==================
    ankles_R = tree_ohs["jointData"][JOINT_MAP["RightAnkle"]]["jointAngle"][:, 2]
    ankles_L = tree_ohs["jointData"][JOINT_MAP["LeftAnkle"]]["jointAngle"][:, 2]

    ankle_R = float(np.max(ankles_R))
    ankle_L = float(np.max(ankles_L))

    # ================== RETURN ==================
    return {
        "sternum_max": sternum_max,
        "pelvis_lat": pelvis_lat,
        "hip_R": hip_R,
        "hip_L": hip_L,
        "knee_R": knee_R,
        "knee_L": knee_L,
        "abd_R": valgo_R,
        "abd_L": valgo_L,
        "ankle_R": ankle_R,
        "ankle_L": ankle_L,
    }

def compute_stepdown(tree, side="R"):
    # ================== TRONCO Y PELVIS ==================
    stern = tree["segmentData"][5]["orientation"]
    pelvis = tree["segmentData"][1]["orientation"]

    eu_s = quat2eul_matlab(stern, "ZYX")
    eu_p = quat2eul_matlab(pelvis, "ZYX")

    sternum_max = float(np.max(eu_s[:, 1]))

    pelvis_max = float(np.max(eu_p[:, 2]))
    pelvis_min = float(np.min(eu_p[:, 2]))

    if pelvis_max <= 5 and pelvis_min >= -5:
        pelvis_lat = "NO"
    else:
        pelvis_lat = (
            pelvis_max
            if abs(pelvis_max) >= abs(pelvis_min)
            else pelvis_min
        )

    # ================== SELECCIÓN DE LADO ==================
    if side == "R":
        hip   = tree["jointData"][JOINT_MAP["RightHip"]]["jointAngle"][:, 2]
        knee  = tree["jointData"][JOINT_MAP["RightKnee"]]["jointAngle"][:, 2]
        abd   = tree["jointData"][JOINT_MAP["RightKnee"]]["jointAngle"][:, 0]
        ankle = tree["jointData"][JOINT_MAP["RightAnkle"]]["jointAngle"][:, 2]
    else:
        hip   = tree["jointData"][JOINT_MAP["LeftHip"]]["jointAngle"][:, 2]
        knee  = tree["jointData"][JOINT_MAP["LeftKnee"]]["jointAngle"][:, 2]
        abd   = tree["jointData"][JOINT_MAP["LeftKnee"]]["jointAngle"][:, 0]
        ankle = tree["jointData"][JOINT_MAP["LeftAnkle"]]["jointAngle"][:, 2]

    hip_max   = float(np.max(hip))
    knee_max  = float(np.max(knee))
    ankle_max = float(np.max(ankle))

    # -------- Valgo rodilla (SOLO MÁXIMO, umbral 5°) --------
    abd_max = float(np.max(abd))
    valgo = abd_max if abd_max > 5 else "NO"

    return {
        "sternum_max": sternum_max,
        "pelvis_lat": pelvis_lat,
        "hip": hip_max,
        "knee": knee_max,
        "abd": valgo,
        "ankle": ankle_max,
    }



# ===================================================================
# 5. CREAR PDF
# ===================================================================


def draw_header_footer(canvas, doc):
    canvas.saveState()

    # Fondo blanco limpio
    canvas.setFillColor(colors.white)
    canvas.rect(0, 0, A4[0], A4[1], stroke=0, fill=1)

    # Barra superior (corporativa)
    canvas.setFillColor(colors.HexColor("#2E5F7F"))  # azul clínico
    canvas.rect(0, A4[1] - 2.2*cm, A4[0], 2.2*cm, stroke=0, fill=1)

    # Título arriba (izquierda)
    canvas.setFillColor(colors.white)
    canvas.setFont("Helvetica-Bold", 12)
    canvas.drawString(2*cm, A4[1] - 1.4*cm, "Informe de Movilidad")

    # Logo (derecha)
    logo_path = LOGO_PORTADA
    try:
        logo = ImageReader(logo_path)
        canvas.drawImage(logo, A4[0] - 5.2*cm, A4[1] - 1.75*cm, width=3.0*cm, height=1.0*cm, mask="auto")
    except:
        pass

    # Pie de página
    canvas.setFillColor(colors.HexColor("#7b8796"))
    canvas.setFont("Helvetica", 9)
    canvas.drawRightString(A4[0] - 2*cm, 1.2*cm, f"Página {doc.page}")

    canvas.restoreState()



def block_tabla_con_imagen(tabla_flowable, imagen_flowable=None):
    """
    Devuelve un bloque con 2 columnas:
    - Izquierda: tabla
    - Derecha: imagen (si existe)
    """
    if imagen_flowable is None:
        return tabla_flowable

    layout = Table(
        [[tabla_flowable, "", imagen_flowable]],
        colWidths=[11*cm, 0.6*cm, 4*cm]
    )
    layout.hAlign = "LEFT"

    layout.setStyle(TableStyle([
        # ✅ Que todo esté arriba (evita que la imagen caiga hacia abajo)
        ("VALIGN", (0, 0), (-1, -1), "TOP"),

        # ✅ Imagen pegada arriba a la derecha
        ("ALIGN", (2, 0), (2, 0), "RIGHT"),
        ("VALIGN", (2, 0), (2, 0), "TOP"),

        # opcional: tabla alineada arriba
        ("VALIGN", (0, 0), (0, 0), "TOP"),

        ("LEFTPADDING", (0, 0), (-1, -1), 0),
        ("RIGHTPADDING", (0, 0), (-1, -1), 0),
        ("TOPPADDING", (0, 0), (-1, -1), 0),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 0),
    ]))

    return layout



def build_pdf(nombre, results, fecha_mvnx):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=2*cm,
        rightMargin=2*cm,
        topMargin=3.2*cm,
        bottomMargin=2*cm
    )


    elements = []

    title_style = ParagraphStyle(
        name="Title",
        fontSize=20,
        leading=24,
        alignment=1,  # centrado
        spaceAfter=20
    )

    section_style = ParagraphStyle(
        name="Section",
        fontSize=13,
        leading=16,
        spaceAfter=8,
        textColor=colors.HexColor("#2E5F7F"),  # azul profesional
        fontName="Helvetica-Bold"
    )



   

    fecha = fecha_mvnx.strftime("%d/%m/%Y")

    portada_title = ParagraphStyle(
        name="PortadaTitle",
        fontSize=26,
        leading=30,
        alignment=1,
        textColor=colors.HexColor("#2E5F7F"),
        fontName="Helvetica-Bold",
    )

    portada_sub = ParagraphStyle(
        name="PortadaSub",
        fontSize=13,
        leading=18,
        alignment=1,
        textColor=colors.HexColor("#4f5b66"),
    )

    elements.append(Spacer(1, 120))
    elements.append(Paragraph("INFORME CLÍNICO DE MOVILIDAD", portada_title))
    elements.append(Spacer(1, 25))
    elements.append(Paragraph(f"<b>Paciente:</b> {nombre}", portada_sub))
    elements.append(Paragraph(f"<b>Fecha:</b> {fecha}", portada_sub))
   
    elements.append(Spacer(1, 250))

    if os.path.exists(LOGO_PORTADA):
        logo = Image(LOGO_PORTADA)
        logo.drawWidth = 10 * cm
        logo.drawHeight = 3.5 * cm
        logo.hAlign = "CENTER"
        elements.append(logo)
    else:
        elements.append(Paragraph("RX2", portada_sub))  # fallback si no encuentra el logo

    elements.append(PageBreak())




    # -------- TABLAS POR PRUEBA --------
    for prueba, filas in results.items():

        # Construimos tabla en formato:
        # | MÉTRICA | DERECHA | IZQUIERDA |
        header = ["", "DERECHA", "IZQUIERDA"]

        tabla = [header]


        # Convertimos la lista en un diccionario agrupado por métrica:
        metrics = {}

        for metric, value in filas:
            m = metric.lower()

            if "derecha" in m or "der." in m:
                key = metric.replace(" derecha", "").replace(" der.", "")
                metrics.setdefault(key, {})["der"] = value

            elif "izquierda" in m or "izq." in m:
                key = metric.replace(" izquierda", "").replace(" izq.", "")
                metrics.setdefault(key, {})["izq"] = value

            else:
                # métrica única (sin lado)
                metrics.setdefault(metric, {})["der"] = value
                metrics.setdefault(metric, {})["izq"] = None

        # Crear filas de tabla
        spans = []

        PASTEL_GREEN = colors.HexColor("#BFE6C3")
        PASTEL_YELLOW = colors.HexColor("#FFF3B0")  # amarillo suave



        # Aplicar estilo
        style = [
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2E5F7F")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 11),

            ("ALIGN", (1, 0), (-1, -1), "CENTER"),
            ("ALIGN", (0, 0), (0, -1), "LEFT"),

            ("LEFTPADDING", (0, 0), (-1, -1), 10),
            ("RIGHTPADDING", (0, 0), (-1, -1), 10),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),

            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F3F6FA")]),

            ("LINEBELOW", (0, 0), (-1, 0), 1, colors.HexColor("#2E5F7F")),
            ("BOX", (0, 0), (-1, -1), 0.6, colors.HexColor("#C9D2DE")),
            ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#C9D2DE")),
        ]




        for i, (metric, vals) in enumerate(metrics.items(), start=1):
            der = vals.get("der")
            izq = vals.get("izq")

            # ---------- MÉTRICA ÚNICA ----------
            if der is not None and izq is None:
                if isinstance(der, str):
                    tabla.append([metric.upper(), der, ""])
                    if der == "NO":
                        style.append(("BACKGROUND", (1, i), (1, i), PASTEL_GREEN))
                else:
                    tabla.append([metric.upper(), f"{der:.1f}°", ""])
                            # ✅ Si es inclinación lateral de pelvis (numérico) -> amarillo si abs>5
                    if "inclinación lateral de la pelvis" in metric.lower():
                        if isinstance(der, (int, float)) and abs(der) > 5:
                            style.append(("BACKGROUND", (1, i), (1, i), PASTEL_YELLOW))
                    



                spans.append(("SPAN", (1, i), (2, i)))

            # ---------- MÉTRICA BILATERAL ----------
            else:
                d = der
                i_ = izq

                # Formateo texto
                d_txt = d if isinstance(d, str) else (f"{d:.1f}°" if d is not None else "")
                i_txt = i_ if isinstance(i_, str) else (f"{i_:.1f}°" if i_ is not None else "")

                tabla.append([metric.upper(), d_txt, i_txt])

                                # 👉 Colorear "NO" en verde
                if d == "NO":
                    style.append(("BACKGROUND", (1, i), (1, i), PASTEL_GREEN))
                if i_ == "NO":
                    style.append(("BACKGROUND", (2, i), (2, i), PASTEL_GREEN))


                # ✅ Pintar en AMARILLO la inclinación lateral si es > 5 (OHS / Step Down)
                if "inclinación lateral de la pelvis" in metric.lower():
                    # derecha
                    if isinstance(d, (int, float)) and abs(d) > 5:
                        style.append(("BACKGROUND", (1, i), (1, i), PASTEL_YELLOW))
                    # izquierda
                    if isinstance(i_, (int, float)) and abs(i_) > 5:
                        style.append(("BACKGROUND", (2, i), (2, i), PASTEL_YELLOW))


                if "dorsiflexión tobillo" in metric.lower():
                    # derecha
                    if isinstance(d, (int, float)) and d < 35:
                        style.append(("BACKGROUND", (1, i), (1, i), PASTEL_YELLOW))
                    # izquierda
                    if isinstance(i_, (int, float)) and i_ < 35:
                        style.append(("BACKGROUND", (2, i), (2, i), PASTEL_YELLOW))                
            
                


        # Crear tabla reportlab
        t = Table(tabla)
        t.hAlign = "LEFT"

    
        # 👇 aplicar las celdas combinadas
        style.extend(spans)

        t.setStyle(TableStyle(style))
                # =========================
        # BLOQUE: TABLA + IMAGEN
        # =========================



                # =========================
        # TAMAÑOS DE IMAGEN (GLOBAL)
        # =========================
        IMG_NORMAL = 3.5 * cm
        IMG_BIG = 4.2 * cm
        IMG_OHS = 3.2 * cm

        if prueba == "Rotación torácica":
            pass


        img = None

        # ----- Caso especial OHS (2 imágenes) -----
        if (
            prueba == "OHS"
            and os.path.exists(IMAGE_MAP["OHS"])
            and os.path.exists(IMAGE_MAP["OHSf"])
        ):
            img_lat = Image(IMAGE_MAP["OHS"])
            img_lat.drawWidth = IMG_OHS
            img_lat.drawHeight = IMG_OHS

            img_front = Image(IMAGE_MAP["OHSf"])
            img_front.drawWidth = IMG_OHS
            img_front.drawHeight = IMG_OHS

            col_img = Table(
                [[img_lat, img_front]],
                colWidths=[IMG_OHS, IMG_OHS]
            )

            col_img.setStyle(TableStyle([
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("LEFTPADDING", (0, 0), (-1, -1), 0),
                ("RIGHTPADDING", (0, 0), (-1, -1), 0),
                ("TOPPADDING", (0, 0), (-1, -1), 0),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 0),
            ]))

            img = col_img


        # ----- Caso Step Down / Rotación de cadera (imagen grande) -----
        elif prueba in ["Step Down", "Rotación de cadera"] and os.path.exists(IMAGE_MAP.get(prueba, "")):
            img = Image(IMAGE_MAP[prueba])

            # ✅ Ajuste especial para Step Down (más alargada)
            if prueba == "Step Down":
                img.drawWidth = 4.0 * cm
                img.drawHeight = 6.0 * cm
            else:
                # Rotación de cadera normal (cuadrada)
                img.drawWidth = IMG_BIG
                img.drawHeight = IMG_BIG

            img.hAlign = "CENTER"



        # ----- Resto de pruebas (imagen normal) -----
        elif prueba in IMAGE_MAP and os.path.exists(IMAGE_MAP[prueba]):
            img = Image(IMAGE_MAP[prueba])
            img.drawWidth = IMG_NORMAL
            img.drawHeight = IMG_NORMAL
            img.hAlign = "CENTER"


        # ✅ OHS y Step Down → FOTO ABAJO
        if prueba in ["OHS", "Step Down"]:
            bloque_test = [
                Paragraph(prueba, section_style),
                Spacer(1, 4),
                t,
                Spacer(1, 8),
            ]

            if img is not None:
                bloque_test.append(img)
                bloque_test.append(Spacer(1, 5))

            elements.append(KeepTogether(bloque_test))

        # ✅ Resto de pruebas → FOTO AL LADO (como lo tenías)
        else:
            bloque_test = [
                Paragraph(prueba, section_style),
                Spacer(1, 4),
                block_tabla_con_imagen(t, img),
                Spacer(1, 5)
            ]

            elements.append(KeepTogether(bloque_test))






    doc.build(
    elements,
    onFirstPage=draw_header_footer,
    onLaterPages=draw_header_footer
    )

    return buffer.getvalue()


def reemplazar_texto_en_presentacion(prs, replacements: dict):
    """
    replacements = {"{{NOMBRE}}": "Mireia", "{{FECHA}}": "20/01/2026"}
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for k, v in replacements.items():
                        if k in run.text:
                            run.text = run.text.replace(k, v)


def get_first_table(slide):
    """
    Devuelve la primera tabla que encuentre en una slide.
    """
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


def agrupar_metricas_bilateral(filas):
    """
    Convierte:
      ("Ángulo pierna derecha", 90), ("Ángulo pierna izquierda", 80)
    en:
      ("Ángulo pierna", 90, 80)
    """
    metrics = {}

    for metric, value in filas:
        m = metric.lower()

        if "derecha" in m or "der." in m:
            key = metric.replace(" derecha", "").replace(" der.", "").strip()
            metrics.setdefault(key, {})["der"] = value

        elif "izquierda" in m or "izq." in m:
            key = metric.replace(" izquierda", "").replace(" izq.", "").strip()
            metrics.setdefault(key, {})["izq"] = value

        else:
            metrics.setdefault(metric, {})["der"] = value
            metrics.setdefault(metric, {})["izq"] = None

    tabla_final = []
    for k, vals in metrics.items():
        tabla_final.append((k, vals.get("der"), vals.get("izq")))

    return tabla_final

def rellenar_tabla_ppt(table, filas_convertidas):
    """
    table: tabla de la plantilla
    filas_convertidas: lista (métrica, derecha, izquierda)
    """

    PASTEL_VERDE = (191, 230, 195)   # "#BFE6C3"
    PASTEL_AMARILLO = (255, 243, 176)  # "#FFF3B0"

    n_cols = len(table.columns)
    n_rows = len(table.rows)

    # ============================================================
    # ✅ CASO 1: TABLA CON 3 COLUMNAS -> MÉTRICA | DERECHA | IZQUIERDA
    # ============================================================
    if n_cols >= 3:
        row = 1  # fila 0 es cabecera

        for metric, der, izq in filas_convertidas:
            if row >= n_rows:
                break

            # METRICA (solo esta columna a tamaño 12)
            set_metric_text_12(table.cell(row, 0), metric.upper())
            set_cell_center(table.cell(row, 0))


            # DERECHA
            if isinstance(der, (int, float)):
                table.cell(row, 1).text = f"{der:.1f}°"
            else:
                table.cell(row, 1).text = "" if der is None else str(der)
            set_cell_center(table.cell(row, 1))

            # IZQUIERDA
            if isinstance(izq, (int, float)):
                table.cell(row, 2).text = f"{izq:.1f}°"
            else:
                table.cell(row, 2).text = "" if izq is None else str(izq)
            set_cell_center(table.cell(row, 2))

            # ==========================================
            # ✅ COLORES SOLO CUANDO APLICA
            # (el resto queda igual que plantilla)
            # ==========================================

            # 👉 Verde si es "NO"
            if der == "NO":
                set_cell_bg(table.cell(row, 1), PASTEL_VERDE)
            if izq == "NO":
                set_cell_bg(table.cell(row, 2), PASTEL_VERDE)

            # 👉 Amarillo si es inclinación lateral pelvis y abs > 5
            if "inclinación lateral de la pelvis" in metric.lower():

                if isinstance(der, (int, float)) and abs(der) > 5:
                    set_cell_bg(table.cell(row, 1), PASTEL_AMARILLO)

                if isinstance(izq, (int, float)) and abs(izq) > 5:
                    set_cell_bg(table.cell(row, 2), PASTEL_AMARILLO)

            # ✅ Amarillo si ankle < 35°
            if "dorsiflexión tobillo" in metric.lower():

                if isinstance(der, (int, float)) and der < 35:
                    set_cell_bg(table.cell(row, 1), PASTEL_AMARILLO)

                if isinstance(izq, (int, float)) and izq < 35:
                    set_cell_bg(table.cell(row, 2), PASTEL_AMARILLO)


            row += 1




        return

    # ============================================================
    # ✅ CASO 2: TABLA CON 2 COLUMNAS -> DERECHA | IZQUIERDA
    # ============================================================
    if n_cols == 2 and n_rows >= 2:

        metric, der, izq = filas_convertidas[0]

        # derecha abajo
        if isinstance(der, (int, float)):
            table.cell(1, 0).text = f"{der:.1f}°"
        else:
            table.cell(1, 0).text = "" if der is None else str(der)

        # izquierda abajo
        if isinstance(izq, (int, float)):
            table.cell(1, 1).text = f"{izq:.1f}°"
        else:
            table.cell(1, 1).text = "" if izq is None else str(izq)

        set_cell_center(table.cell(1, 0))
        set_cell_center(table.cell(1, 1))

        # ✅ Solo pintamos cuando aplica
        if der == "NO":
            set_cell_bg(table.cell(1, 0), PASTEL_VERDE)
        if izq == "NO":
            set_cell_bg(table.cell(1, 1), PASTEL_VERDE)

        return


def build_pptx(nombre, results, fecha_mvnx, tipo_biomecanica="Simple"):


    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

    TEMPLATE_PATH_SIMPLE   = os.path.join(BASE_DIR, "PLANTILLA_MOVILIDAD.pptx")
    TEMPLATE_PATH_CARRERA  = os.path.join(BASE_DIR, "PLANTILLA_CARRERA.pptx")
    TEMPLATE_PATH_COMPLETA = os.path.join(BASE_DIR, "PLANTILLA_COMPLETA.pptx")



    if tipo_biomecanica == "Carrera":
        template_to_use = TEMPLATE_PATH_CARRERA
    elif tipo_biomecanica == "Completa":
        template_to_use = TEMPLATE_PATH_COMPLETA
    else:
        template_to_use = TEMPLATE_PATH_SIMPLE


    prs = Presentation(template_to_use)



    fecha = fecha_mvnx.strftime("%d/%m/%Y")

    # ✅ 1) Cambiar nombre y fecha en la plantilla (portada)
    reemplazar_texto_en_presentacion(prs, {
        "{{NOMBRE}}": nombre,
        "{{FECHA}}": fecha
    })

    # ✅ 2) MAPA: cada prueba -> qué slide es en tu plantilla
    # ⚠️ esto lo tienes que ajustar a tu PPTX real
    SLIDE_MAP = {
        "ASLR": 2,
        "AKE": 3,
        "Rotación de cadera": 4,
        "Rotación torácica": 5,
        "Thomas": 6,
        "Lunge": 7,
        "Flexión plantar": 8,
        "OHS": 9,
        "Step Down": 10,
    }

    # ✅ 1) Rellenar solo los test que existen
    for prueba, filas in results.items():

        if prueba not in SLIDE_MAP:
            continue

        slide_index = SLIDE_MAP[prueba]
        if slide_index >= len(prs.slides):
            continue

        slide = prs.slides[slide_index]
        table = get_first_table(slide)
        if table is None:
            continue

        filas_convertidas = agrupar_metricas_bilateral(filas)
        rellenar_tabla_ppt(table, filas_convertidas)

    # ✅ 2) BORRAR diapositivas de pruebas NO realizadas
    pruebas_realizadas = set(results.keys())

    slides_to_delete = []
    for prueba, idx in SLIDE_MAP.items():
        if prueba not in pruebas_realizadas:
            slides_to_delete.append(idx)

    # 🔥 borrar en orden inverso para que no se rompan los índices
    for idx in sorted(slides_to_delete, reverse=True):
        if idx < len(prs.slides):
            delete_slide(prs, idx)

    # ✅ Exportar pptx
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()

def safe_load(path):
    return load_mvnx_struct(path) if os.path.exists(path) else None

# ===================================================================
# 6. STREAMLIT APP
# ===================================================================

TEST_CATALOG = [
    ("ASLR", "ASLR"),
    ("AKE", "AKE"),
    ("Rotación de cadera", "ROT_CADERA"),
    ("Rotación torácica", "ROT_TORACICA"),
    ("Thomas derecha", "THOMAS_D"),
    ("Thomas izquierda", "THOMAS_I"),
    ("Lunge", "LUNGE"),
    ("Flexión plantar", "PLANTAR"),
    ("OHS", "OHS"),
    ("Step Down derecha", "STEPDOWN_D"),
    ("Step Down izquierda", "STEPDOWN_I"),
]



def main():
    st.title("🦵 Informe de movilidad MVN (.mvnx)")
    nombre = st.text_input("Nombre del paciente")
    zipfile_u = st.file_uploader("Sube ZIP con los .mvnx", type=["zip"])
    st.subheader("Selecciona los ejercicios realizados")

    selected_tests = []

    for label, key in TEST_CATALOG:
        if st.checkbox(label, value=True):
            selected_tests.append(key)

    st.subheader("Sistema XSens")
    xsens_mode = st.radio(
        "¿Qué sistema se ha usado?",
        ["Awinda", "Link (traje)"],
        horizontal=True
    )


    st.subheader("Biomecánica de carrera")
    bio_carrera = st.radio(
        "Tipo de biomecánica",
        ["Simple", "Carrera", "Completa"],
        horizontal=True
    )

    tipo_biomecanica = bio_carrera   # "Simple" | "Carrera" | "Completa"



    formato = st.radio("Formato de informe:", ["PDF", "PowerPoint"], horizontal=True)

    if st.button("Generar informe"):

        if not nombre or not zipfile_u:
            st.error("Falta nombre o archivo ZIP.")
            return
        
        with tempfile.TemporaryDirectory() as tmp:
            with zipfile.ZipFile(zipfile_u, "r") as z:
                z.extractall(tmp)

            # cargar archivos
     

            # 1️⃣ listar los mvnx reales del ZIP
            

            mvnx_files = []

            for root, _, files in os.walk(tmp):
                for f in files:
                    if f.lower().endswith(".mvnx") and not f.startswith("._"):
                        mvnx_files.append(os.path.join(root, f))

            # 🔒 filtrar solo MVNX que sean XML reales
            mvnx_files = [f for f in mvnx_files if is_valid_mvnx(f)]

            if not mvnx_files:
                st.error("❌ No se encontraron archivos MVNX válidos en el ZIP.")
                return

            mvnx_files.sort()

            fecha_mvnx = get_fecha_from_mvnx(mvnx_files[0])

            if fecha_mvnx is None:
                st.error("❌ No se ha podido leer la fecha del MVNX.")
                return






            # 2️⃣ asignarlos por orden a los tests seleccionados
            loaded_tests = {}

            file_idx = 0  # índice REAL sobre los archivos

            for _, key in TEST_CATALOG:
                if key not in selected_tests:
                    continue  # test no realizado → no consume archivo

                if file_idx >= len(mvnx_files):
                    break  # no hay más archivos

                # path = os.path.join(tmp, mvnx_files[file_idx])
                path = mvnx_files[file_idx]
                loaded_tests[key] = load_mvnx_struct(path)

                file_idx += 1            

            results = {}

            # ================= ASLR =================
            if loaded_tests.get("ASLR") is not None:
                R_aslr, L_aslr = compute_aslr(loaded_tests["ASLR"])
                results["ASLR"] = [
                    ("Ángulo pierna derecha", R_aslr),
                    ("Ángulo pierna izquierda", L_aslr),
                ]

            # ================= AKE =================
            if loaded_tests.get("AKE") is not None:
                R_ake, L_ake, R_frame, L_frame = compute_ake(loaded_tests["AKE"])
                results["AKE"] = [
                    ("Flexión rodilla derecha", R_ake),
                    ("Flexión rodilla izquierda", L_ake),
                ]

            # ================= ROTACIÓN CADERA =================
            if loaded_tests.get("ROT_CADERA") is not None:
                R_int, R_ext, L_int, L_ext = compute_rot_cadera(
                    loaded_tests["ROT_CADERA"],
                    xsens_mode
                )

                R_rom = R_int + R_ext
                L_rom = L_int + L_ext

                results["Rotación de cadera"] = [
                    ("Rot. interna derecha", R_int),
                    ("Rot. externa derecha", R_ext),
                    ("ROM total derecha", R_rom),
                    ("Rot. interna izquierda", L_int),
                    ("Rot. externa izquierda", L_ext),
                    ("ROM total izquierda", L_rom),
                ]



            # ================= ROTACIÓN TORÁCICA =================
            if loaded_tests.get("ROT_TORACICA") is not None:

                # OJO: ahora tree incluye jointAngleErgo
                tree_tor = loaded_tests["ROT_TORACICA"]

                # Crear estructura compatible con compute_rot_toracica
                # porque compute_rot_toracica espera tree["jointAngleErgo"]
                tree_for_tor = {"jointAngleErgo": tree_tor["jointAngleErgo"]}

                R_torso, L_torso = compute_rot_toracica(tree_for_tor)

                results["Rotación torácica"] = [
                    ("Rotación torácica derecha", R_torso),
                    ("Rotación torácica izquierda", L_torso),
                ]



            # ================= THOMAS =================
            if (
                loaded_tests.get("THOMAS_D") is not None
                and loaded_tests.get("THOMAS_I") is not None
            ):
                th_R_thigh, th_R_knee, th_L_thigh, th_L_knee = compute_thomas(
                    loaded_tests["THOMAS_D"],
                    loaded_tests["THOMAS_I"],
                    xsens_mode
                )
                results["Thomas"] = [
                    ("Orientación muslo derecha", th_R_thigh),
                    ("Orientación muslo izquierda", th_L_thigh),
                    ("Flexión rodilla derecha", th_R_knee),
                    ("Flexión rodilla izquierda", th_L_knee),
                ]


            # ================= LUNGE =================
            if loaded_tests.get("LUNGE") is not None:
                lunge_R, lunge_L = compute_lunge(loaded_tests["LUNGE"])
                results["Lunge"] = [
                    ("Dorsiflexión tobillo derecha", lunge_R),
                    ("Dorsiflexión tobillo izquierda", lunge_L),
                ]


            # ================= FLEXIÓN PLANTAR =================
            if loaded_tests.get("PLANTAR") is not None:
                pf_R, pf_L = compute_plantarflex(loaded_tests["PLANTAR"])
                results["Flexión plantar"] = [
                    ("Flexión plantar tobillo derecha", pf_R),
                    ("Flexión plantar tobillo izquierda", pf_L),
                ]

            # ================= OHS =================
            if loaded_tests.get("OHS") is not None:
                ohs = compute_ohs(loaded_tests["OHS"])
                results["OHS"] = [
                    ("Inclinación anterior del tronco", ohs["sternum_max"]),
                    ("Inclinación lateral de la pelvis", ohs["pelvis_lat"]),
                    ("Flexión cadera derecha", ohs["hip_R"]),
                    ("Flexión cadera izquierda", ohs["hip_L"]),
                    ("Flexión rodilla derecha", ohs["knee_R"]),
                    ("Flexión rodilla izquierda", ohs["knee_L"]),
                    ("Valgo rodilla derecha", ohs["abd_R"]),
                    ("Valgo rodilla izquierda", ohs["abd_L"]),
                    ("Dorsiflexión tobillo derecha", ohs["ankle_R"]),
                    ("Dorsiflexión tobillo izquierda", ohs["ankle_L"]),
                ]

            # ================= STEP DOWN =================
            if (
                loaded_tests.get("STEPDOWN_D") is not None
                and loaded_tests.get("STEPDOWN_I") is not None
            ):
                sd_R = compute_stepdown(loaded_tests["STEPDOWN_D"], "R")
                sd_L = compute_stepdown(loaded_tests["STEPDOWN_I"], "L")

                results["Step Down"] = [
                    ("Inclinación anterior del tronco derecha", sd_R["sternum_max"]),
                    ("Inclinación anterior del tronco izquierda", sd_L["sternum_max"]),
                    ("Inclinación lateral de la pelvis derecha", sd_R["pelvis_lat"]),
                    ("Inclinación lateral de la pelvis izquierda", sd_L["pelvis_lat"]),
                    ("Flexión cadera derecha", sd_R["hip"]),
                    ("Flexión cadera izquierda", sd_L["hip"]),
                    ("Flexión rodilla derecha", sd_R["knee"]),
                    ("Flexión rodilla izquierda", sd_L["knee"]),
                    ("Valgo rodilla derecha", sd_R["abd"]),
                    ("Valgo rodilla izquierda", sd_L["abd"]),
                    ("Dorsiflexión tobillo derecha", sd_R["ankle"]),
                    ("Dorsiflexión tobillo izquierda", sd_L["ankle"]),
                ]

            if formato == "PDF":
                pdf = build_pdf(nombre, results, fecha_mvnx)

                st.success("✅ Informe PDF generado correctamente.")
                st.download_button(
                    "📄 Descargar PDF",
                    pdf,
                    f"Informe_{nombre}.pdf",
                    mime="application/pdf",
                )

            else:
                pptx = build_pptx(nombre, results, fecha_mvnx, tipo_biomecanica=tipo_biomecanica)



                st.success("✅ Informe PowerPoint generado correctamente.")
                st.download_button(
                    "📊 Descargar PowerPoint",
                    pptx,
                    f"Informe_{nombre}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )



if __name__ == "__main__":
    main()